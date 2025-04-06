#!/usr/bin/env python3
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk
import re
import os
import socket
import json
import logging
import hashlib
import threading
import queue
from typing import List, Dict, Tuple, Optional, Pattern, Callable, Any, Set
from dataclasses import dataclass, field
import time
import functools
import secrets

# Set up logging with rotation
from logging.handlers import RotatingFileHandler

# Create logs directory if it doesn't exist
os.makedirs('logs', exist_ok=True)
log_handler = RotatingFileHandler(
    filename='logs/sanitizer.log',
    maxBytes=5 * 1024 * 1024,  # 5 MB
    backupCount=3
)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[log_handler]
)

# Caching decorator for pattern compilation
def cached_regex_compile(func):
    """Cache compiled regex patterns to improve performance"""
    _cache = {}
    
    @functools.wraps(func)
    def wrapper(pattern, flags=0):
        # Create a unique cache key considering both pattern and flags
        cache_key = (pattern, flags)
        
        if cache_key not in _cache:
            _cache[cache_key] = func(pattern, flags)
        
        return _cache[cache_key]
    
    return wrapper

# Wrap re.compile with caching
compile_regex = cached_regex_compile(re.compile)

@dataclass
class SanitizationPattern:
    """Enhanced data class to hold sanitization pattern details"""
    name: str
    pattern: Pattern
    replacement: str
    hash_flag: bool
    enabled: bool = True
    case_sensitive: bool = False
    comment: str = ""
    
    def __post_init__(self):
        """Validate pattern on initialization"""
        if not self.name:
            raise ValueError("Pattern name cannot be empty")
        
    def apply(self, text: str) -> Tuple[str, int]:
        """Apply the pattern and return the modified text and count of replacements"""
        if not self.enabled:
            return text, 0
            
        matches = self.pattern.findall(text)
        count = len(matches)
        
        if count > 0:
            def replace_match(match):
                matched_text = match.group(0)
                if self.hash_flag:
                    # More secure hashing with salt
                    salt = secrets.token_hex(4)
                    return f"<HASHED_{self.name.upper()}_{hashlib.sha256((salt + matched_text).encode('utf-8')).hexdigest()[:12].upper()}>"
                return self.replacement if match.lastindex is None else match.expand(self.replacement)
            
            result = self.pattern.sub(replace_match, text)
            return result, count
        
        return text, 0


class SanitizerEngine:
    """Enhanced core sanitization engine"""
    
    def __init__(self):
        self.patterns: Dict[str, SanitizationPattern] = {}
        self.username = os.getlogin()
        self.hostname = socket.gethostname()
        self.initialize_default_patterns()
        
    def initialize_default_patterns(self):
        """Set up the default sanitization patterns with improved error handling"""
        default_patterns = {
            "Email": (r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b", "<EMAIL>", False),
            "IP": (r"\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b", "<IP>", False),
            "GUID": (r"\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b", "<GUID>", False),
            "SID": (r"\bS-1-5-21-\d+-\d+-\d+-\d+\b", "<SID>", True),
            "DistinguishedName_CN": (r"(?<=distinguishedName=CN=)[^,]+", "<USER>", False),
            "SAMAccountName": (r"SAMAccountName=([^,\s]+)", "SAMAccountName=<USER>", False),
            "DisplayName": (r"DisplayName=([^,\n]+)", "DisplayName=<USER>", False),
            "Timestamp": (r"(?<=Logged at )\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}", "<TIMESTAMP>", False),
            "LegacyExchangeDN": (r"legacyExchangeDN=([^,\n]+)", "legacyExchangeDN=<LEGACY_DN>", False),
            "MailNickname": (r"mailNickname=([^,\s]+)", "mailNickname=<NICKNAME>", False),
            "ObjectCategory": (r"objectCategory=([^,\n]+)", "objectCategory=<CATEGORY>", False),
            "ProxyAddresses": (r"proxyAddresses=([^,\n]+)", "proxyAddresses=<PROXY>", False),
            "ShowInAddressBook": (r"showInAddressBook=([^,\n]+)", "showInAddressBook=<ADDRESS_BOOK>", False)
        }
        
        try:
            for name, (pattern, replacement, hash_flag) in default_patterns.items():
                compiled_pattern = compile_regex(pattern)
                self.patterns[name] = SanitizationPattern(
                    name, 
                    compiled_pattern, 
                    replacement, 
                    hash_flag
                )
            
            # System-specific patterns with improved detection
            system_patterns = [
                ("Username", self.username, r"\b{}\b", "<USERNAME>"),
                ("Hostname", self.hostname, r"\b{}\b", "<HOSTNAME>"),
                ("FilePath_Username", self.username, r"(?<=Users\\){}", "<USERNAME>")
            ]
            
            for name, value, template, replacement in system_patterns:
                try:
                    # Escape special regex characters in username/hostname
                    escaped_value = re.escape(value)
                    pattern = compile_regex(template.format(escaped_value), re.IGNORECASE)
                    self.patterns[name] = SanitizationPattern(
                        name, 
                        pattern, 
                        replacement, 
                        False, 
                        case_sensitive=False
                    )
                except Exception as e:
                    logging.warning(f"Could not create system pattern {name}: {e}")
            
        except Exception as e:
            logging.error(f"Failed to compile default patterns: {e}")
            raise ValueError(f"Pattern initialization failed: {e}")
    
    def load_config(self, config_path: str) -> Tuple[bool, str]:
        """Enhanced config loading with more robust validation"""
        if not os.path.exists(config_path):
            return False, f"Config file not found: {config_path}"
            
        try:
            with open(config_path, 'r', encoding='utf-8') as file:
                config = json.load(file)
                custom_patterns = config.get("patterns", [])
            
            # Comprehensive pattern validation
            validated_patterns = []
            for pattern in custom_patterns:
                # Check required fields
                if not all(key in pattern for key in ["name", "regex", "replacement"]):
                    logging.warning(f"Skipping invalid pattern: missing required fields {pattern}")
                    continue
                
                name = pattern.get("name", "Unnamed")
                regex = pattern.get("regex")
                replacement = pattern.get("replacement")
                
                # Validate regex
                try:
                    flags = 0 if pattern.get("case_sensitive", False) else re.IGNORECASE
                    compiled_regex = compile_regex(regex, flags)
                    
                    validated_patterns.append(SanitizationPattern(
                        name=name,
                        pattern=compiled_regex,
                        replacement=replacement,
                        hash_flag=pattern.get("hash", False),
                        case_sensitive=pattern.get("case_sensitive", False),
                        comment=pattern.get("comment", "")
                    ))
                except re.error as e:
                    logging.warning(f"Invalid regex for pattern '{name}': {e}")
            
            # Update patterns
            self.patterns.update({p.name: p for p in validated_patterns})
            
            logging.info(f"Successfully loaded {len(validated_patterns)} custom patterns from {config_path}")
            return True, f"Loaded {len(validated_patterns)} custom patterns"
                
        except json.JSONDecodeError as e:
            return False, f"Invalid JSON in config file: {e}"
        except Exception as e:
            logging.error(f"Failed to load config {config_path}: {e}")
            return False, f"Error loading config: {e}"
    
    def sanitize(self, log_text: str, callback: Optional[Callable[[str, int, int], None]] = None) -> Tuple[str, Dict[str, int]]:
        """Enhanced sanitization with improved performance and tracing"""
        sanitized_text = log_text
        pattern_counts = {}
        total_patterns = len(self.patterns)
        
        # Create a copy of patterns to avoid modification during iteration
        patterns_copy = list(self.patterns.items())
        
        for i, (name, pattern_obj) in enumerate(patterns_copy):
            try:
                sanitized_text, count = pattern_obj.apply(sanitized_text)
                
                # Only add to counts if there were replacements
                if count > 0:
                    pattern_counts[name] = count
                
                # Optional progress callback
                if callback and i % 3 == 0:
                    callback(sanitized_text, i + 1, total_patterns)
                    
            except Exception as e:
                logging.error(f"Error applying pattern '{name}': {e}")
        
        return sanitized_text, pattern_counts



class LogSanitizerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Advanced Log Sanitizer")
        self.root.minsize(800, 600)
        
        # Set up the engine
        self.engine = SanitizerEngine()
        
        # Status variables
        self.status_var = tk.StringVar(value="Ready")
        self.progress_var = tk.IntVar(value=0)
        
        # Current office file path
        self.current_office_file = None
        
        # Task queue for background operations
        self.task_queue = queue.Queue()
        
        # Set up UI
        self.setup_gui()
        
        # Try to load default config
        self.load_default_config()
        
        # Start background worker
        self.start_background_worker()
    
    def setup_gui(self):
        """Initialize the GUI components"""
        # Configure root window grid
        self.root.grid_rowconfigure(0, weight=1)  # Main content expands
        self.root.grid_rowconfigure(1, weight=0)  # Status bar fixed height
        self.root.grid_rowconfigure(2, weight=0)  # Sanitize button fixed height
        self.root.grid_columnconfigure(0, weight=1)  # Full width
        
        # Create a main frame with padding
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky="nsew")
        
        # Make main_frame expandable
        main_frame.grid_rowconfigure(0, weight=1)
        main_frame.grid_columnconfigure(0, weight=1)
        
        # Create horizontal paned window for left panel and main content
        h_pane = ttk.PanedWindow(main_frame, orient=tk.HORIZONTAL)
        h_pane.grid(row=0, column=0, sticky="nsew")
        
        # Left panel for pattern configuration
        left_panel = ttk.Frame(h_pane, padding="5")
        h_pane.add(left_panel, weight=1)
        
        # Main content area
        content_frame = ttk.Frame(h_pane, padding="5")
        h_pane.add(content_frame, weight=3)
        
        # Set up the left panel (patterns)
        self.setup_pattern_panel(left_panel)
        
        # Set up the content area (input/output)
        self.setup_content_area(content_frame)
        
        # Status bar at the bottom
        status_frame = ttk.Frame(self.root)
        status_frame.grid(row=1, column=0, sticky="ew", padx=10, pady=(0, 5))
        
        ttk.Label(status_frame, textvariable=self.status_var).pack(side=tk.LEFT)
        ttk.Progressbar(status_frame, variable=self.progress_var, length=200, mode="determinate").pack(side=tk.RIGHT)
        
        # Add sanitize button at the bottom - always visible
        sanitize_btn = ttk.Button(self.root, text="SANITIZE LOG", command=self.sanitize, style="Accent.TButton")
        sanitize_btn.grid(row=2, column=0, sticky="ew", padx=10, pady=(0, 10))
        
        # Style the button to make it stand out
        self.root.style = ttk.Style()
        self.root.style.configure("Accent.TButton", font=("", 12, "bold"))
    
    def setup_pattern_panel(self, parent):
        """Set up the pattern configuration panel"""
        # Configure parent grid
        parent.grid_rowconfigure(0, weight=0)  # Label fixed height
        parent.grid_rowconfigure(1, weight=1)  # Pattern tree expands
        parent.grid_rowconfigure(2, weight=0)  # Button frame fixed height
        parent.grid_rowconfigure(3, weight=0)  # Config frame fixed height
        parent.grid_columnconfigure(0, weight=1)  # Full width
        
        # Label for the section
        ttk.Label(parent, text="Sanitization Patterns", font=("", 10, "bold")).grid(row=0, column=0, sticky="w")
        
        # Frame for pattern list
        pattern_frame = ttk.Frame(parent)
        pattern_frame.grid(row=1, column=0, sticky="nsew", pady=(5, 0))
        
        # Configure pattern frame grid
        pattern_frame.grid_rowconfigure(0, weight=1)
        pattern_frame.grid_columnconfigure(0, weight=1)
        pattern_frame.grid_columnconfigure(1, weight=0)
        
        # Scrollbars
        vsb = ttk.Scrollbar(pattern_frame, orient=tk.VERTICAL)
        hsb = ttk.Scrollbar(pattern_frame, orient=tk.HORIZONTAL)
        
        # Pattern treeview
        self.pattern_tree = ttk.Treeview(
            pattern_frame, 
            columns=("status",), 
            height=10,
            yscrollcommand=vsb.set,
            xscrollcommand=hsb.set
        )
        
        # Configure scrollbars
        vsb.config(command=self.pattern_tree.yview)
        hsb.config(command=self.pattern_tree.xview)
        
        # Configure treeview
        self.pattern_tree.heading("#0", text="Pattern Name")
        self.pattern_tree.heading("status", text="Enabled")
        self.pattern_tree.column("#0", width=150)
        self.pattern_tree.column("status", width=50, anchor=tk.CENTER)
        
        # Bind double-click to toggle pattern
        self.pattern_tree.bind("<Double-1>", lambda event: self.toggle_pattern())
        
        # Place components
        self.pattern_tree.grid(row=0, column=0, sticky="nsew")
        vsb.grid(row=0, column=1, sticky="ns")
        hsb.grid(row=1, column=0, sticky="ew")
        
        # Populate patterns
        self.refresh_pattern_list()
        
        # Pattern control buttons
        btn_frame = ttk.Frame(parent)
        btn_frame.grid(row=2, column=0, sticky="ew", pady=(5, 0))
        
        ttk.Button(btn_frame, text="Add", command=self.add_pattern).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(btn_frame, text="Edit", command=self.edit_pattern).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(btn_frame, text="Enable/Disable", command=self.toggle_pattern).pack(side=tk.LEFT)
        
        # Config file section
        config_frame = ttk.LabelFrame(parent, text="Configuration")
        config_frame.grid(row=3, column=0, sticky="ew", pady=(10, 0))
        
        # Configure config frame grid
        config_frame.grid_rowconfigure(0, weight=0)
        config_frame.grid_rowconfigure(1, weight=0)
        config_frame.grid_columnconfigure(0, weight=1)
        
        ttk.Button(config_frame, text="Load Config", command=self.load_config).grid(row=0, column=0, sticky="ew", pady=(5, 0))
        ttk.Button(config_frame, text="Save Config", command=self.save_config).grid(row=1, column=0, sticky="ew", pady=(5, 5))
    
    def setup_content_area(self, parent):
        """Set up the main content area with input/output panels"""
        # Configure parent grid
        parent.grid_rowconfigure(0, weight=1)  # Input/output panes expand
        parent.grid_columnconfigure(0, weight=1)  # Full width
        
        # Create vertical paned window for input and output
        v_pane = ttk.PanedWindow(parent, orient=tk.VERTICAL)
        v_pane.grid(row=0, column=0, sticky="nsew")
        
        # Input panel
        input_frame = ttk.LabelFrame(v_pane, text="Log Input")
        v_pane.add(input_frame, weight=1)
        
        # Configure input frame grid
        input_frame.grid_rowconfigure(0, weight=1)  # Text area expands
        input_frame.grid_rowconfigure(1, weight=0)  # Button frame fixed height
        input_frame.grid_columnconfigure(0, weight=1)  # Full width
        
        # Input text area with line numbers
        self.input_text = scrolledtext.ScrolledText(input_frame, wrap=tk.WORD, undo=True)
        self.input_text.grid(row=0, column=0, sticky="nsew", padx=5, pady=5)
        
        # Add keyboard bindings for paste
        self.input_text.bind("<Control-v>", self.paste_from_clipboard_event)
        
        # Input control buttons
        input_btn_frame = ttk.Frame(input_frame)
        input_btn_frame.grid(row=1, column=0, sticky="ew", padx=5, pady=(0, 5))
        
        ttk.Button(input_btn_frame, text="Load from File", command=self.load_file).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(input_btn_frame, text="Load Word/Office File", command=self.load_office_file).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(input_btn_frame, text="Paste from Clipboard", command=self.paste_from_clipboard).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(input_btn_frame, text="Clear", command=lambda: self.input_text.delete(1.0, tk.END)).pack(side=tk.LEFT)
        
        # Output panel
        output_frame = ttk.LabelFrame(v_pane, text="Sanitized Output")
        v_pane.add(output_frame, weight=1)
        
        # Configure output frame grid
        output_frame.grid_rowconfigure(0, weight=1)  # Text area expands
        output_frame.grid_rowconfigure(1, weight=0)  # Button frame fixed height
        output_frame.grid_columnconfigure(0, weight=1)  # Full width
        
        # Output text area
        self.output_text = scrolledtext.ScrolledText(output_frame, wrap=tk.WORD)
        self.output_text.grid(row=0, column=0, sticky="nsew", padx=5, pady=5)
        
        # Output control buttons
        output_btn_frame = ttk.Frame(output_frame)
        output_btn_frame.grid(row=1, column=0, sticky="ew", padx=5, pady=(0, 5))
        
        ttk.Button(output_btn_frame, text="Save to File", command=self.save_file).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(output_btn_frame, text="Save as Word File", command=self.save_as_office_file).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(output_btn_frame, text="Copy to Clipboard", command=self.copy_to_clipboard).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(output_btn_frame, text="Clear", command=lambda: self.output_text.delete(1.0, tk.END)).pack(side=tk.LEFT)
    
    def refresh_pattern_list(self):
        """Refresh the pattern treeview"""
        # Clear existing items
        for item in self.pattern_tree.get_children():
            self.pattern_tree.delete(item)
            
        # Add patterns
        for name, pattern in self.engine.patterns.items():
            status = "✓" if pattern.enabled else "✗"
            self.pattern_tree.insert("", tk.END, text=name, values=(status,))
    
    def add_pattern(self):
        """Add a new sanitization pattern"""
        # Create dialog
        dialog = tk.Toplevel(self.root)
        dialog.title("Add Pattern")
        dialog.transient(self.root)
        dialog.resizable(False, False)
        dialog.grab_set()
        
        # Set up form
        ttk.Label(dialog, text="Pattern Name:").grid(row=0, column=0, sticky=tk.W, padx=5, pady=5)
        name_var = tk.StringVar(dialog)
        ttk.Entry(dialog, textvariable=name_var, width=30).grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(dialog, text="Regular Expression:").grid(row=1, column=0, sticky=tk.W, padx=5, pady=5)
        regex_var = tk.StringVar(dialog)
        ttk.Entry(dialog, textvariable=regex_var, width=30).grid(row=1, column=1, padx=5, pady=5)
        
        ttk.Label(dialog, text="Replacement:").grid(row=2, column=0, sticky=tk.W, padx=5, pady=5)
        replacement_var = tk.StringVar(dialog)
        ttk.Entry(dialog, textvariable=replacement_var, width=30).grid(row=2, column=1, padx=5, pady=5)
        
        # Options
        case_sensitive_var = tk.BooleanVar(dialog, value=False)
        ttk.Checkbutton(dialog, text="Case Sensitive", variable=case_sensitive_var).grid(row=3, column=0, sticky=tk.W, padx=5, pady=5)
        
        hash_var = tk.BooleanVar(dialog, value=False)
        ttk.Checkbutton(dialog, text="Hash Values", variable=hash_var).grid(row=3, column=1, sticky=tk.W, padx=5, pady=5)
        
        # Status message
        status_var = tk.StringVar(dialog)
        ttk.Label(dialog, textvariable=status_var, foreground="red").grid(row=4, column=0, columnspan=2, padx=5, pady=5)
        
        # Buttons
        btn_frame = ttk.Frame(dialog)
        btn_frame.grid(row=5, column=0, columnspan=2, padx=5, pady=5)
        
        def validate_and_add():
            name = name_var.get().strip()
            regex = regex_var.get().strip()
            replacement = replacement_var.get()
            
            if not name:
                status_var.set("Pattern name is required")
                return
                
            if not regex:
                status_var.set("Regular expression is required")
                return
                
            if name in self.engine.patterns:
                status_var.set(f"Pattern '{name}' already exists")
                return
            
            try:
                flags = 0 if case_sensitive_var.get() else re.IGNORECASE
                compiled_pattern = re.compile(regex, flags)
                
                self.engine.patterns[name] = SanitizationPattern(
                    name,
                    compiled_pattern,
                    replacement,
                    hash_var.get()
                )
                
                self.refresh_pattern_list()
                dialog.destroy()
                
            except re.error as e:
                status_var.set(f"Invalid regex: {e}")
        
        ttk.Button(btn_frame, text="Add", command=validate_and_add).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(btn_frame, text="Cancel", command=dialog.destroy).pack(side=tk.LEFT)
        
        # Center dialog
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (self.root.winfo_width() // 2) - (width // 2) + self.root.winfo_x()
        y = (self.root.winfo_height() // 2) - (height // 2) + self.root.winfo_y()
        dialog.geometry(f"{width}x{height}+{x}+{y}")
    
    def edit_pattern(self):
        """Edit the selected pattern"""
        selected = self.pattern_tree.selection()
        if not selected:
            messagebox.showinfo("Info", "Please select a pattern to edit")
            return
            
        pattern_name = self.pattern_tree.item(selected[0], "text")
        if pattern_name not in self.engine.patterns:
            messagebox.showerror("Error", f"Pattern '{pattern_name}' not found")
            return
            
        pattern = self.engine.patterns[pattern_name]
        
        # Create dialog
        dialog = tk.Toplevel(self.root)
        dialog.title(f"Edit Pattern: {pattern_name}")
        dialog.transient(self.root)
        dialog.resizable(False, False)
        dialog.grab_set()
        
        # Set up form
        ttk.Label(dialog, text="Regular Expression:").grid(row=0, column=0, sticky=tk.W, padx=5, pady=5)
        regex_var = tk.StringVar(dialog, value=pattern.pattern.pattern)
        ttk.Entry(dialog, textvariable=regex_var, width=30).grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(dialog, text="Replacement:").grid(row=1, column=0, sticky=tk.W, padx=5, pady=5)
        replacement_var = tk.StringVar(dialog, value=pattern.replacement)
        ttk.Entry(dialog, textvariable=replacement_var, width=30).grid(row=1, column=1, padx=5, pady=5)
        
        # Options
        flags = pattern.pattern.flags
        case_sensitive_var = tk.BooleanVar(dialog, value=(flags & re.IGNORECASE) == 0)
        ttk.Checkbutton(dialog, text="Case Sensitive", variable=case_sensitive_var).grid(row=2, column=0, sticky=tk.W, padx=5, pady=5)
        
        hash_var = tk.BooleanVar(dialog, value=pattern.hash_flag)
        ttk.Checkbutton(dialog, text="Hash Values", variable=hash_var).grid(row=2, column=1, sticky=tk.W, padx=5, pady=5)
        
        enabled_var = tk.BooleanVar(dialog, value=pattern.enabled)
        ttk.Checkbutton(dialog, text="Enabled", variable=enabled_var).grid(row=3, column=0, sticky=tk.W, padx=5, pady=5)
        
        # Status message
        status_var = tk.StringVar(dialog)
        ttk.Label(dialog, textvariable=status_var, foreground="red").grid(row=4, column=0, columnspan=2, padx=5, pady=5)
        
        # Buttons
        btn_frame = ttk.Frame(dialog)
        btn_frame.grid(row=5, column=0, columnspan=2, padx=5, pady=5)
        
        def validate_and_update():
            regex = regex_var.get().strip()
            replacement = replacement_var.get()
            
            if not regex:
                status_var.set("Regular expression is required")
                return
            
            try:
                flags = 0 if case_sensitive_var.get() else re.IGNORECASE
                compiled_pattern = re.compile(regex, flags)
                
                self.engine.patterns[pattern_name] = SanitizationPattern(
                    pattern_name,
                    compiled_pattern,
                    replacement,
                    hash_var.get(),
                    enabled_var.get()
                )
                
                self.refresh_pattern_list()
                dialog.destroy()
                
            except re.error as e:
                status_var.set(f"Invalid regex: {e}")
        
        ttk.Button(btn_frame, text="Update", command=validate_and_update).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(btn_frame, text="Cancel", command=dialog.destroy).pack(side=tk.LEFT)
        
        # Center dialog
        dialog.update_idletasks()
        width = dialog.winfo_width()
        height = dialog.winfo_height()
        x = (self.root.winfo_width() // 2) - (width // 2) + self.root.winfo_x()
        y = (self.root.winfo_height() // 2) - (height // 2) + self.root.winfo_y()
        dialog.geometry(f"{width}x{height}+{x}+{y}")
    
    def toggle_pattern(self):
        """Toggle the selected pattern's enabled status"""
        selected = self.pattern_tree.selection()
        if not selected:
            messagebox.showinfo("Info", "Please select a pattern to toggle")
            return
            
        pattern_name = self.pattern_tree.item(selected[0], "text")
        if pattern_name not in self.engine.patterns:
            messagebox.showerror("Error", f"Pattern '{pattern_name}' not found")
            return
            
        pattern = self.engine.patterns[pattern_name]
        pattern.enabled = not pattern.enabled
        
        self.refresh_pattern_list()
    
    def load_default_config(self):
        """Load the default configuration file if available"""
        default_config_path = os.path.join(os.path.dirname(__file__), "config.json")
        if os.path.exists(default_config_path):
            success, message = self.engine.load_config(default_config_path)
            if success:
                self.status_var.set(f"Loaded default config: {os.path.basename(default_config_path)}")
                self.refresh_pattern_list()
            else:
                logging.warning(f"Could not load default config: {message}")
    
    def load_config(self):
        """Load a configuration file"""
        file_path = filedialog.askopenfilename(
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")],
            title="Load Configuration"
        )
        
        if file_path:
            # Add task to queue
            self.task_queue.put(("load_config", file_path))
    
    def save_config(self):
        """Save the current configuration to a file"""
        file_path = filedialog.asksaveasfilename(
            defaultextension=".json",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")],
            title="Save Configuration"
        )
        
        if file_path:
            # Prepare config data
            config = {
                "patterns": []
            }
            
            for name, pattern in self.engine.patterns.items():
                config["patterns"].append({
                    "name": name,
                    "regex": pattern.pattern.pattern,
                    "replacement": pattern.replacement,
                    "case_sensitive": (pattern.pattern.flags & re.IGNORECASE) == 0,
                    "hash": pattern.hash_flag,
                    "enabled": pattern.enabled
                })
            
            try:
                with open(file_path, 'w', encoding='utf-8') as file:
                    json.dump(config, file, indent=2)
                
                self.status_var.set(f"Config saved: {os.path.basename(file_path)}")
                logging.info(f"Saved config to: {file_path}")
            except Exception as e:
                logging.error(f"Failed to save config {file_path}: {e}")
                messagebox.showerror("Error", f"Failed to save config: {e}")
                self.status_var.set("Error saving config")
    
    def load_file(self):
        """Load a log file for sanitization"""
        file_path = filedialog.askopenfilename(
            filetypes=[("Text files", "*.txt *.log"), ("All files", "*.*")],
            title="Load Log File"
        )
        
        if file_path:
            # Add task to queue
            self.task_queue.put(("load_file", file_path))
    
    def save_file(self):
        """Save the sanitized output to a file"""
        file_path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("Log files", "*.log"), ("All files", "*.*")],
            title="Save Sanitized Log"
        )
        
        if file_path:
            try:
                with open(file_path, 'w', encoding='utf-8') as file:
                    file.write(self.output_text.get(1.0, tk.END))
                
                self.status_var.set(f"Saved to: {os.path.basename(file_path)}")
                logging.info(f"Saved sanitized log to: {file_path}")
            except Exception as e:
                logging.error(f"Failed to save file {file_path}: {e}")
                messagebox.showerror("Error", f"Failed to save file: {e}")
                self.status_var.set("Error saving file")
    
    def paste_from_clipboard_event(self, event=None):
        """Handle Ctrl+V keyboard event for pasting clipboard content"""
        try:
            clipboard_text = self.root.clipboard_get()
            if clipboard_text:
                # Get current selection or cursor position
                try:
                    sel_start = self.input_text.index("sel.first")
                    sel_end = self.input_text.index("sel.last")
                    self.input_text.delete(sel_start, sel_end)
                    self.input_text.insert(sel_start, clipboard_text)
                except tk.TclError:  # No selection
                    self.input_text.insert(tk.INSERT, clipboard_text)
            return "break"  # Prevent default paste behavior
        except Exception as e:
            logging.error(f"Paste from clipboard error: {e}")
            self.status_var.set("Error pasting from clipboard")
            return "break"
    
    def paste_from_clipboard(self):
        """Paste text from clipboard to input area"""
        try:
            clipboard_text = self.root.clipboard_get()
            self.input_text.delete(1.0, tk.END)
            self.input_text.insert(tk.END, clipboard_text)
            self.status_var.set("Pasted from clipboard")
        except Exception as e:
            self.status_var.set("No text in clipboard or error pasting")
    
    def copy_to_clipboard(self):
        """Copy sanitized output to clipboard"""
        output_text = self.output_text.get(1.0, tk.END)
        if output_text.strip():
            self.root.clipboard_clear()
            self.root.clipboard_append(output_text)
            self.status_var.set("Copied to clipboard")
        else:
            self.status_var.set("No output to copy")
            
    def load_office_file(self):
        """Load a Word or other Office file for sanitization"""
        file_path = filedialog.askopenfilename(
            filetypes=[
                ("Word Documents", "*.docx *.doc"),
                ("Excel Workbooks", "*.xlsx *.xls"),
                ("PowerPoint Presentations", "*.pptx *.ppt"),
                ("Text files", "*.txt"),
                ("All files", "*.*")
            ],
            title="Load Office Document"
        )
        
        if file_path:
            # Add task to queue
            self.task_queue.put(("load_office_file", file_path))
    
    def process_load_office_file(self, file_path):
        """Process the load office file task"""
        self.update_status(f"Loading {os.path.basename(file_path)}...")
        self.update_progress(0)
        
        try:
            # Store the current office file path
            self.current_office_file = file_path
            
            file_ext = os.path.splitext(file_path)[1].lower()
            text_content = ""
            
            # Process based on file type
            if file_ext in ['.docx', '.doc']:
                # Handle Word documents
                import docx
                try:
                    doc = docx.Document(file_path)
                    text_content = "\n\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    # Alternative method if python-docx is not available
                    try:
                        from subprocess import Popen, PIPE
                        import tempfile
                        
                        # Try to use a text extraction tool like antiword
                        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as temp_file:
                            temp_path = temp_file.name
                        
                        cmd = ["antiword", file_path]
                        with Popen(cmd, stdout=PIPE, stderr=PIPE) as p:
                            stdout, stderr = p.communicate()
                            if p.returncode == 0:
                                text_content = stdout.decode('utf-8', errors='replace')
                            else:
                                # Fallback to a simple message if extraction fails
                                text_content = f"[Word document: {os.path.basename(file_path)}]\n\nUnable to extract text content directly. Install python-docx or antiword for Word document support."
                    except:
                        text_content = f"[Word document: {os.path.basename(file_path)}]\n\nUnable to extract text. Install python-docx for Word document support."
            
            elif file_ext in ['.xlsx', '.xls']:
                # Handle Excel files
                try:
                    import pandas as pd
                    # Read all sheets and convert to text
                    excel_file = pd.ExcelFile(file_path)
                    sheets_text = []
                    
                    for sheet_name in excel_file.sheet_names:
                        df = excel_file.parse(sheet_name)
                        sheets_text.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
                    
                    text_content = "\n\n".join(sheets_text)
                except ImportError:
                    text_content = f"[Excel workbook: {os.path.basename(file_path)}]\n\nUnable to extract content. Install pandas for Excel support."
            
            elif file_ext in ['.pptx', '.ppt']:
                try:
                    import pptx
                    prs = pptx.Presentation(file_path)
                    presentation_text = []
                    
                    for i, slide in enumerate(prs.slides):
                        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                        presentation_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
                    
                    text_content = "\n\n".join(presentation_text)
                except ImportError:
                    text_content = f"[PowerPoint presentation: {os.path.basename(file_path)}]\n\nUnable to extract content. Install python-pptx for PowerPoint support."
            
            else:
                # Handle as regular text file
                with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                    text_content = file.read()
            
            # Update UI
            self.root.after(0, lambda: self.input_text.delete(1.0, tk.END))
            self.root.after(0, lambda: self.input_text.insert(tk.END, text_content))
            
            self.update_status(f"Loaded Office file: {os.path.basename(file_path)}")
            self.update_progress(100)
            logging.info(f"Loaded office file: {file_path}")
            
        except Exception as e:
            self.update_status(f"Error loading Office file: {e}")
            self.update_progress(0)
            logging.error(f"Failed to load Office file {file_path}: {e}")
            self.root.after(0, lambda: messagebox.showerror("Error", f"Failed to load Office file: {e}"))
    
    def save_as_office_file(self):
        """Save the sanitized content back to an office file format"""
        if not self.output_text.get(1.0, tk.END).strip():
            messagebox.showinfo("Info", "No sanitized content to save")
            return
        
        # Determine default file type based on original file
        default_ext = ".txt"
        filetypes = [
            ("Text files", "*.txt"),
            ("Word Documents", "*.docx"),
            ("All files", "*.*")
        ]
        
        if self.current_office_file:
            # Determine original file type and set appropriate default
            original_ext = os.path.splitext(self.current_office_file)[1].lower()
            original_name = os.path.basename(self.current_office_file)
            base_name = os.path.splitext(original_name)[0]
            
            if original_ext in ['.docx', '.doc']:
                default_ext = ".docx"
                
            # Set default filename with _Sanitized suffix
            default_filename = f"{base_name}_Sanitized{default_ext}"
        else:
            default_filename = "Sanitized_Document.txt"
        
        # Ask user where to save the file
        file_path = filedialog.asksaveasfilename(
            defaultextension=default_ext,
            filetypes=filetypes,
            initialfile=default_filename,
            title="Save Sanitized Document"
        )
        
        if not file_path:
            return
            
        try:
            sanitized_text = self.output_text.get(1.0, tk.END)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.docx':
                # Save as Word document
                try:
                    import docx
                    doc = docx.Document()
                    
                    # Split by paragraph and add to document
                    paragraphs = sanitized_text.split('\n')
                    for para_text in paragraphs:
                        if para_text.strip():  # Skip empty paragraphs
                            doc.add_paragraph(para_text)
                    
                    doc.save(file_path)
                except ImportError:
                    messagebox.showerror("Error", "Python-docx library is required to save Word documents. Please install it using pip: pip install python-docx")
                    return
            else:
                # Save as text file
                with open(file_path, 'w', encoding='utf-8') as file:
                    file.write(sanitized_text)
            
            self.status_var.set(f"Saved to: {os.path.basename(file_path)}")
            logging.info(f"Saved sanitized document to: {file_path}")
            
        except Exception as e:
            logging.error(f"Failed to save document {file_path}: {e}")
            messagebox.showerror("Error", f"Failed to save document: {e}")
            self.status_var.set("Error saving document")
    
    def sanitize(self):
        """Sanitize the log text"""
        log_text = self.input_text.get(1.0, tk.END)
        if not log_text.strip():
            messagebox.showinfo("Info", "No input text to sanitize")
            return
        
        # Add task to queue
        self.task_queue.put(("sanitize", log_text))
    
    def start_background_worker(self):
        """Start the background worker thread"""
        self.worker_running = True
        self.worker_thread = threading.Thread(target=self.background_worker, daemon=True)
        self.worker_thread.start()
    
    def background_worker(self):
        """Background worker that processes tasks from the queue"""
        while self.worker_running:
            try:
                # Get task with a timeout to allow checking worker_running
                try:
                    task = self.task_queue.get(timeout=0.5)
                except queue.Empty:
                    continue
                
                task_type, task_data = task
                
                if task_type == "load_file":
                    self.process_load_file(task_data)
                elif task_type == "load_office_file":
                    self.process_load_office_file(task_data)
                elif task_type == "load_config":
                    self.process_load_config(task_data)
                elif task_type == "sanitize":
                    self.process_sanitize(task_data)
                
                self.task_queue.task_done()
            except Exception as e:
                logging.error(f"Background worker error: {e}")
                self.update_status(f"Error: {e}")
    
    def process_load_file(self, file_path):
        """Process the load file task"""
        self.update_status(f"Loading {os.path.basename(file_path)}...")
        self.update_progress(0)
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                file_content = file.read()
            
            # Update UI
            self.root.after(0, lambda: self.input_text.delete(1.0, tk.END))
            self.root.after(0, lambda: self.input_text.insert(tk.END, file_content))
            
            self.update_status(f"Loaded: {os.path.basename(file_path)}")
            self.update_progress(100)
            logging.info(f"Loaded file: {file_path}")
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    file_content = file.read()
                
                # Update UI
                self.root.after(0, lambda: self.input_text.delete(1.0, tk.END))
                self.root.after(0, lambda: self.input_text.insert(tk.END, file_content))
                
                self.update_status(f"Loaded with latin-1 encoding: {os.path.basename(file_path)}")
                self.update_progress(100)
                logging.info(f"Loaded file with latin-1 encoding: {file_path}")
            except Exception as e:
                self.update_status(f"Error loading file: {e}")
                self.update_progress(0)
                logging.error(f"Failed to load file {file_path}: {e}")
                self.root.after(0, lambda: messagebox.showerror("Error", f"Failed to load file: {e}"))
        except Exception as e:
            self.update_status(f"Error loading file: {e}")
            self.update_progress(0)
            logging.error(f"Failed to load file {file_path}: {e}")
            self.root.after(0, lambda: messagebox.showerror("Error", f"Failed to load file: {e}"))
    
    def process_load_config(self, file_path):
        """Process the load config task"""
        self.update_status(f"Loading config {os.path.basename(file_path)}...")
        self.update_progress(50)
        
        success, message = self.engine.load_config(file_path)
        
        if success:
            self.update_status(f"Config loaded: {os.path.basename(file_path)}")
            self.root.after(0, self.refresh_pattern_list)
        else:
            self.update_status(f"Error loading config: {message}")
            self.root.after(0, lambda: messagebox.showerror("Error", message))
        
        self.update_progress(100)
    
    def process_sanitize(self, log_text):
        """Process the sanitization task"""
        self.update_status("Sanitizing...")
        self.update_progress(0)
        
        start_time = time.time()
        
        # Define the progress callback
        def update_sanitize_progress(current_text, current, total):
            progress = int((current / total) * 100)
            self.update_progress(progress)
            if progress % 25 == 0:  # Update preview every 25%
                self.root.after(0, lambda: self.output_text.delete(1.0, tk.END))
                self.root.after(0, lambda: self.output_text.insert(tk.END, current_text))
        
        # Perform sanitization
        try:
            sanitized_text, pattern_counts = self.engine.sanitize(log_text, update_sanitize_progress)
            
            # Calculate statistics
            elapsed_time = time.time() - start_time
            total_replacements = sum(pattern_counts.values())
            replacements_desc = ", ".join(f"{name}: {count}" for name, count in pattern_counts.items() if count > 0)
            
            # Update UI
            self.root.after(0, lambda: self.output_text.delete(1.0, tk.END))
            self.root.after(0, lambda: self.output_text.insert(tk.END, sanitized_text))
            
            # Show summary
            summary = f"Sanitized in {elapsed_time:.2f}s: {total_replacements} replacements"
            if replacements_desc:
                # Show details in a message box if there are replacements
                self.root.after(0, lambda: messagebox.showinfo("Sanitization Complete", 
                    f"Time: {elapsed_time:.2f} seconds\n"
                    f"Total replacements: {total_replacements}\n\n"
                    f"Pattern details:\n{replacements_desc}"))
            
            self.update_status(summary)
            self.update_progress(100)
            logging.info(f"Sanitization completed: {total_replacements} replacements in {elapsed_time:.2f}s")
        except Exception as e:
            self.update_status(f"Error sanitizing: {e}")
            self.update_progress(0)
            logging.error(f"Sanitization failed: {e}")
            self.root.after(0, lambda: messagebox.showerror("Error", f"Sanitization failed: {e}"))
    
    def update_status(self, message):
        """Update the status bar message"""
        self.root.after(0, lambda: self.status_var.set(message))
    
    def update_progress(self, value):
        """Update the progress bar value"""
        self.root.after(0, lambda: self.progress_var.set(value))

def main():
    # Set up exception handling
    def show_error(*args):
        if len(args) == 4:
            # Called by Tkinter directly with (self, exc_type, exc_value, exc_traceback)
            _, exc_type, exc_value, exc_traceback = args
        else:
            # Called from our code with (exc_type, exc_value, exc_traceback)
            exc_type, exc_value, exc_traceback = args
            
        error_msg = f"An unexpected error occurred:\n{exc_type.__name__}: {exc_value}"
        logging.error("Uncaught exception", exc_info=(exc_type, exc_value, exc_traceback))
        messagebox.showerror("Error", error_msg)
    
    # Create the root window
    root = tk.Tk()
    root.title("Advanced Log Sanitizer")
    root.geometry("1000x700")
    
    # Set minimum window size to ensure elements are visible
    root.minsize(800, 600)
    
    # Set the hook for error handling
    root.report_callback_exception = show_error
    
    # Add required Python packages for Office file support
    required_packages = [
        "python-docx",  # For Word documents
        "pandas",       # For Excel files
        "openpyxl",     # For Excel backend
        "python-pptx"   # For PowerPoint
    ]
    
    # Check for available packages (but don't fail if missing)
    try:
        import importlib.util
        for package in required_packages:
            try:
                spec = importlib.util.find_spec(package.replace("-", ""))
                if spec is None:
                    logging.warning(f"Package {package} not found. Some features may be unavailable.")
            except ImportError:
                logging.warning(f"Could not check for package {package}.")
    except:
        pass
    
    # Create and run the app
    app = LogSanitizerApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()